import json
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 读取JSON文件
with open("analysis_result/tmp.json", "r", encoding="utf-8") as f:
    data = json.load(f)

# 创建PPT
prs = Presentation()


def create_chart_slide(prs, title_text, chart_data, is_sentiment=False):
    # 创建新的幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白布局

    # 设置标题
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True

    # 设置图表位置和大小（居中且更大）
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)

    # 添加柱状图
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # 设置图表样式
    plot = chart.plots[0]
    plot.has_data_labels = True  # 显示数据标签

    # 设置数据标签格式
    data_labels = plot.data_labels
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = RGBColor(0, 0, 0)

    # 设置系列颜色
    series_list = plot.series
    series_list[0].format.fill.solid()
    series_list[0].format.fill.fore_color.rgb = RGBColor(65, 105, 225)  # 深蓝色
    series_list[1].format.fill.solid()
    series_list[1].format.fill.fore_color.rgb = RGBColor(255, 99, 71)  # 珊瑚红

    # 设置图例
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(12)

    # 设置坐标轴
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = Pt(10)
    if len(chart_data.categories) > 8:
        category_axis.tick_labels.rotation = -45

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.font.size = Pt(10)
    if is_sentiment:
        value_axis.maximum_scale = 100.0


# 遍历每个主题
for topic in data.keys():
    # 处理Buzz数据
    buzz_data = data[topic]["pkBuzz"]
    chart_data = CategoryChartData()
    chart_data.categories = [item["name"] for item in buzz_data]
    chart_data.add_series("惠庭", [item["惠庭"] for item in buzz_data])
    chart_data.add_series(
        "行业平均",
        [item.get("industryAverage", item.get("行业平均", 0)) for item in buzz_data],
    )

    # 创建Buzz图表幻灯片
    create_chart_slide(prs, f"{topic} - Buzz对比", chart_data)

    # 处理SentimentScore数据
    sentiment_data = data[topic]["pkSentimentScore"]
    chart_data = CategoryChartData()
    chart_data.categories = [item["name"] for item in sentiment_data]
    chart_data.add_series("惠庭", [float(item["惠庭"]) for item in sentiment_data])
    chart_data.add_series(
        "行业平均",
        [
            float(item.get("industryAverage", item.get("行业平均", 0)))
            for item in sentiment_data
        ],
    )

    # 创建SentimentScore图表幻灯片
    create_chart_slide(prs, f"{topic} - Sentiment Score对比", chart_data, True)

# 保存PPT
prs.save("analysis_result/home2_vs_industry_average.pptx")
